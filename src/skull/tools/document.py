"""Text extraction for common document formats (PDF, DOCX, XLSX, PPTX) and
OCR for images (PNG, JPG, etc.), so read_file/sandbox_read_file can return
readable text instead of raw binary garbage when pointed at one of these.
Dispatches purely on file extension - no content sniffing.

Takes raw bytes rather than a path, since the sandbox variant needs to fetch
bytes over the network before extracting; the local variant just reads the
file first. Keeps one extraction implementation shared by both callers.

Image OCR (pytesseract) is a text-recognition tool, not vision - it reads
text that appears in an image (a screenshot, a scanned page, a photo of a
document). The chat model itself has no image understanding: this endpoint
accepts an image_url content block without erroring, but silently ignores
the pixel data and the model just guesses from the text prompt (confirmed
by testing - it named the wrong color for solid-color test images). OCR is
the only real way to get anything useful out of an image file here.
"""

import io
import shutil

DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
SUPPORTED_EXTENSIONS = DOCUMENT_EXTENSIONS | IMAGE_EXTENSIONS


def is_extractable(path: str) -> bool:
    return _ext(path) in SUPPORTED_EXTENSIONS


def _ext(path: str) -> str:
    dot = path.rfind(".")
    return path[dot:].lower() if dot != -1 else ""


def extract_text(path: str, data: bytes) -> str:
    """Extract readable text from `data` (the raw file bytes), using
    `path`'s extension to pick the right parser. Raises ValueError for an
    unsupported extension - callers should check is_extractable() first."""
    ext = _ext(path)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".docx":
        return _extract_docx(data)
    if ext in (".xlsx", ".xlsm"):
        return _extract_xlsx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext in IMAGE_EXTENSIONS:
        return _extract_image(data)
    raise ValueError(f"unsupported extension for text extraction: {ext!r}")


def ocr_available() -> bool:
    """Whether the tesseract binary is actually installed - pytesseract is
    just a wrapper around it and raises at call time (not import time) if
    it's missing, so callers should check this before advertising OCR as
    available."""
    return shutil.which("tesseract") is not None


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        pages.append(f"--- page {i} ---\n{text}" if text else f"--- page {i} (no extractable text) ---")
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join("" if c is None else str(c) for c in row))
        sheets.append(f"--- sheet: {ws.title} ---\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
        slides.append(f"--- slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_image(data: bytes) -> str:
    if not ocr_available():
        raise RuntimeError(
            "OCR requires the tesseract binary, which isn't installed "
            "(install it with your system package manager, e.g. "
            "`brew install tesseract` on macOS)"
        )

    import pytesseract
    from PIL import Image

    image = Image.open(io.BytesIO(data))
    text = pytesseract.image_to_string(image).strip()
    return text if text else "(no text detected in image)"
