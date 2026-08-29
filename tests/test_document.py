"""Tests for tools/document.py - text extraction for PDF/DOCX/XLSX/PPTX.
Each format is round-tripped: build a real minimal file in-memory with the
same library used to author documents elsewhere in the app, then confirm
extract_text() recovers the content."""

import io

import pytest

from skull.tools import document


def test_is_extractable_recognizes_supported_extensions():
    assert document.is_extractable("resume.pdf")
    assert document.is_extractable("resume.PDF")  # case-insensitive
    assert document.is_extractable("report.docx")
    assert document.is_extractable("data.xlsx")
    assert document.is_extractable("data.xlsm")
    assert document.is_extractable("slides.pptx")


def test_is_extractable_recognizes_image_extensions():
    assert document.is_extractable("screenshot.png")
    assert document.is_extractable("photo.jpg")
    assert document.is_extractable("photo.JPEG")
    assert document.is_extractable("scan.tiff")


def test_is_extractable_rejects_other_extensions():
    assert not document.is_extractable("notes.txt")
    assert not document.is_extractable("script.py")
    assert not document.is_extractable("no_extension")


def test_extract_text_raises_for_unsupported_extension():
    with pytest.raises(ValueError):
        document.extract_text("notes.txt", b"hello")


def test_extract_pdf_returns_page_text():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)

    text = document.extract_text("doc.pdf", buf.getvalue())
    # A blank page has no extractable text, but must still be reported per-page
    # rather than raising or returning nothing.
    assert "page 1" in text


def test_extract_docx_returns_paragraph_text():
    import docx

    doc = docx.Document()
    doc.add_paragraph("Hello from a test document.")
    doc.add_paragraph("Second paragraph.")
    buf = io.BytesIO()
    doc.save(buf)

    text = document.extract_text("doc.docx", buf.getvalue())
    assert "Hello from a test document." in text
    assert "Second paragraph." in text


def test_extract_docx_includes_table_content():
    import docx

    doc = docx.Document()
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Name"
    table.rows[0].cells[1].text = "Score"
    buf = io.BytesIO()
    doc.save(buf)

    text = document.extract_text("doc.docx", buf.getvalue())
    assert "Name" in text
    assert "Score" in text


def test_extract_xlsx_returns_cell_values():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    buf = io.BytesIO()
    wb.save(buf)

    text = document.extract_text("data.xlsx", buf.getvalue())
    assert "Sheet1" in text
    assert "Name" in text
    assert "Alice" in text
    assert "30" in text


def test_extract_pptx_returns_slide_text():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "My Slide Title"
    buf = io.BytesIO()
    prs.save(buf)

    text = document.extract_text("slides.pptx", buf.getvalue())
    assert "slide 1" in text
    assert "My Slide Title" in text


def test_extract_pdf_malformed_data_raises():
    with pytest.raises(Exception):
        document.extract_text("doc.pdf", b"not a real pdf")


# ---------------------------------------------------------------------------
# Image OCR - skipped entirely if tesseract isn't installed on this machine,
# since pytesseract is just a wrapper around the external binary.
# ---------------------------------------------------------------------------

requires_tesseract = pytest.mark.skipif(
    not document.ocr_available(), reason="tesseract binary not installed"
)


def _render_text_image(text: str) -> bytes:
    import io as _io

    from PIL import Image, ImageDraw

    img = Image.new("RGB", (400, 100), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), text, fill="black")
    buf = _io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@requires_tesseract
def test_extract_image_ocr_reads_rendered_text():
    data = _render_text_image("HELLO WORLD")
    text = document.extract_text("screenshot.png", data)
    assert "HELLO" in text.upper()


@requires_tesseract
def test_extract_image_with_no_text_returns_placeholder():
    from PIL import Image

    img = Image.new("RGB", (50, 50), color="white")
    buf = io.BytesIO()
    img.save(buf, format="PNG")

    text = document.extract_text("blank.png", buf.getvalue())
    assert text == "(no text detected in image)"


def test_extract_image_raises_clear_error_when_tesseract_missing(monkeypatch):
    monkeypatch.setattr(document, "ocr_available", lambda: False)
    with pytest.raises(RuntimeError, match="tesseract"):
        document.extract_text("photo.png", b"fake image bytes")
